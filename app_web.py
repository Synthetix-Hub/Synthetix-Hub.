import streamlit as st
import pytesseract
from PIL import Image
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io
from docx import Document # WICHTIG: Hier oben laden!

# ======================================================
# SYNTHETIX HUB - KONFIGURATION
# ======================================================
GROQ_API_KEY = st.secrets["GROQ_API_KEY"]
PEXELS_API_KEY = st.secrets["PEXELS_API_KEY"]

st.set_page_config(page_title="Synthetix Hub | AI OS", page_icon="🌌", layout="wide")

# Custom CSS
st.markdown("""
    <style>
    .main-title { font-size: 45px; font-weight: 800; color: #00ffa6; margin-bottom: 10px; }
    .stButton>button { background-color: #00ffa6 !important; color: black !important; font-weight: bold !important; border-radius: 10px !important; }
    </style>
    """, unsafe_allow_html=True)

# --- SIDEBAR ---
with st.sidebar:
    st.markdown("## 🌌 Synthetix Hub")
    st.divider()
    menu = st.radio("Tool auswählen:", ["📽️ Präsentations-Generator", "📝 Text-Optimierer", "🎓 Homework Engine", "⚙️ Einstellungen"])
    st.divider()
    st.caption("Version 1.2.1 | Admin Mode")

# --- TOOL 1: PRÄSENTATIONEN ---
# --- TOOL 1: PRÄSENTATIONEN ---
if menu == "📽️ Präsentations-Generator":
    st.markdown('<p class="main-title">🎓 Präsentations-Generator</p>', unsafe_allow_html=True)
    t_in = st.text_input("Thema der Präsentation:")
    
    if st.button("🚀 Engine starten") and t_in:
        with st.spinner("Synthetix Hub arbeitet..."):
            # 1. KI fragt Groq nach Inhalten
            headers = {"Authorization": f"Bearer {GROQ_API_KEY}"}
            payload = {
                "model": "llama-3.3-70b-versatile",
                "messages": [{"role": "user", "content": f"Erstelle 5 Folien über {t_in}. Gib nur den Text aus, getrennt durch '---'."}]
            }
            res = res = requests.post("https://api.groq.com/openai/v1/chat/completions", json=payload, headers=headers)
            
            if res.status_code == 200:
                result_json = res.json()
                content = result_json['choices'][0]['message']['content']
                
                # Hier geht es weiter mit # 2. PowerPoint Datei erstellen
                prs = Presentation()
                # ... (dein restlicher Code für die Folien)
            else:
                st.error(f"Fehler von Groq: {res.status_code}")
                st.write(res.text) # Das zeigt uns genau, WARUM der Key abgelehnt wird
            content = res.json()['choices'][0]['message']['content']

            # 2. PowerPoint Datei erstellen
            prs = Presentation()
            for slide_text in content.split('---'):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = t_in
                slide.placeholders[1].text = slide_text.strip()

            # 3. Datei im Speicher bereitstellen
            pptx_io = io.BytesIO()
            prs.save(pptx_io)
            pptx_io.seek(0)

            st.success("Präsentation fertig!")
            st.download_button(label="📥 Datei herunterladen", data=pptx_io, file_name="Praesentation.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
# --- TOOL 2: TEXT-OPTIMIERER ---
elif menu == "📝 Text-Optimierer":
    st.markdown('<p class="main-title">📝 Text-Optimierer</p>', unsafe_allow_html=True)
    raw = st.text_area("Dein Text:", height=200, placeholder="Tippe hier...")
    if st.button("✨ Veredeln") and raw:
        st.info("Poliere Text...")

# --- TOOL 3: HOMEWORK ENGINE ---
elif menu == "🎓 Homework Engine":
    st.markdown('<p class="main-title">🎓 Homework Engine</p>', unsafe_allow_html=True)
    col1, col2 = st.columns(2)
    with col1:
        aufgabe = st.text_area("Aufgabe/Zusatzinfos:", height=150, placeholder="Was soll ich tun?")
        name_input = st.text_input("Dein Name:", "Schüler")
    with col2:
        img_file = st.file_uploader("Arbeitsblatt scannen", type=["jpg", "png", "jpeg"])

    if st.button("🔍 Lösung generieren & Word erstellen"):
        if aufgabe or img_file:
            with st.spinner("Synthetix arbeitet..."):
                input_data = aufgabe
                if img_file:
                    img = Image.open(img_file)
                    input_data = pytesseract.image_to_string(img, lang='deu') + "\n" + aufgabe
                
                # API Call
                headers = {"Authorization": f"Bearer {GROQ_API_KEY}"}
                payload = {"model": "llama-3.3-70b-versatile", "messages": [{"role": "user", "content": f"Löse ausführlich: {input_data}"}]}
                res = requests.post("https://api.groq.com/openai/v1/chat/completions", json=payload, headers=headers).json()
                loesung = res['choices'][0]['message']['content']
                
                # WORD BAUEN
                doc = Document()
                doc.add_heading('Hausaufgaben-Lösung', 0)
                doc.add_paragraph(f"Name: {name_input}")
                doc.add_paragraph(loesung)
                doc_io = io.BytesIO()
                doc.save(doc_io)
                doc_io.seek(0)
                
                st.markdown("### ✅ Lösung gefunden!")
                st.write(loesung)
                st.download_button("📥 Als Word (.docx) herunterladen", doc_io, f"Hausaufgabe_{name_input}.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")

# --- SETTINGS ---
elif menu == "⚙️ Einstellungen":
    st.markdown('<p class="main-title">⚙️ Einstellungen</p>', unsafe_allow_html=True)
    st.write("System: Online")
